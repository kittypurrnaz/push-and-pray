from __future__ import annotations

import anthropic
import json
import io
import re
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Palette ────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy  — cover backgrounds
INDIGO    = RGBColor(0x63, 0x66, 0xF1)   # indigo     — primary accent
INDIGO_D  = RGBColor(0x49, 0x4B, 0xC8)   # dark indigo — end slide bottom bar
INDIGO_P  = RGBColor(0xEE, 0xF2, 0xFF)   # pale indigo — card backgrounds
SLATE     = RGBColor(0x1E, 0x29, 0x3B)   # dark slate — headings on white
MID       = RGBColor(0x47, 0x55, 0x69)   # mid grey   — body text
LIGHT     = RGBColor(0x94, 0xA3, 0xB8)   # light grey — captions / footer
SOFT      = RGBColor(0xC7, 0xD2, 0xFE)   # soft indigo — end slide secondary text
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)   # subtle card borders
SHADOW    = RGBColor(0xD1, 0xD5, 0xDB)   # card drop-shadow colour

H1  = "Calibri Light"   # large display headings
H2  = "Calibri"         # section headings in band
TXT = "Calibri"         # body text

SW = 13.33   # slide width  (inches, 16:9 wide)
SH = 7.5     # slide height (inches)

SYSTEM_PROMPT = (
    "You are an expert post-event report writer for a professional events agency. "
    "You write clear, data-driven, client-facing reports. Always respond with valid JSON only — "
    "no markdown fences, no commentary outside the JSON object."
)


class ReportGenerator:

    def __init__(self, api_key: str):
        self.client = anthropic.Anthropic(api_key=api_key)

    # ── Prompt assembly ────────────────────────────────────────────────────────

    def _build_user_prompt(
        self,
        config: dict,
        attendance_summary: dict | None,
        survey_summary: dict | None,
        benchmark_context: dict | None = None,
        event_context: dict | None = None,
        analyst_findings: dict | None = None,
        memory_context: dict | None = None,
    ) -> str:
        sections = config.get("sections", {})

        lines = [
            "Generate a post-event report based on the details and data below.\n",
            "EVENT DETAILS:",
            f"  Event Name : {config['event_name']}",
            f"  Event Date : {config['event_date']}",
            f"  Client     : {config.get('client_name') or 'Not specified'}",
            f"  Organizer  : {config.get('organizer_name') or 'Not specified'}",
            f"  Context    : {config.get('additional_context') or 'None provided'}",
            "",
        ]

        if attendance_summary:
            lines.append("ATTENDANCE DATA (cleaned summary):")
            lines.append(json.dumps(attendance_summary, indent=2, default=str))
            lines.append("")

        if survey_summary:
            lines.append("SURVEY / FEEDBACK DATA (cleaned summary):")
            lines.append(json.dumps(survey_summary, indent=2, default=str))
            lines.append("")

        if benchmark_context and benchmark_context.get("found") and benchmark_context.get("context"):
            lines.append(f"INDUSTRY BENCHMARKS (sourced from {benchmark_context.get('source', 'Genspark')}):")
            lines.append(benchmark_context["context"])
            lines.append(
                "Use these benchmarks to contextualise the client's numbers — "
                "reference them explicitly where relevant."
            )
            lines.append("")

        if event_context and event_context.get("found") and event_context.get("context"):
            lines.append(f"EVENT AUDIENCE CONTEXT (sourced from {event_context.get('source', 'Genspark')}):")
            lines.append(event_context["context"])
            lines.append(
                "Use this to frame recommendations and highlights in terms of what this "
                "audience specifically values — make the report feel researched, not templated."
            )
            lines.append("")

        if analyst_findings:
            lines.append("RICH ANALYTICS (computed from cleaned data):")
            lines.append(json.dumps(analyst_findings, indent=2, default=str))
            lines.append(
                "These are statistically computed findings. Prefer these numbers over the "
                "raw summary above when they conflict — they are more accurate."
            )
            lines.append("")

        if memory_context and memory_context.get("total_events_processed", 0) > 0:
            lines.append("YOUR HISTORICAL CONTEXT (from past events processed by this system):")
            lines.append(json.dumps(memory_context, indent=2, default=str))
            lines.append(
                "Where relevant, compare this event's performance against the client's own "
                "historical averages (not just industry benchmarks). "
                "Note improvements or regressions."
            )
            lines.append("")

        lines.append("Return a JSON object with these keys (include only the keys listed):")
        lines.append("{")

        if sections.get("exec_summary"):
            lines.append(
                '  "executive_summary": "3-4 sentence paragraph summarising the event success, '
                'key metrics, and overall outcome for the client",'
            )
        if sections.get("attendance"):
            lines.append(
                '  "attendance_analysis": "Paragraph on registration vs attendance numbers, '
                'show-up rate, and any audience composition insights. Quote specific figures '
                'and compare to industry benchmarks where available.",'
            )
        if sections.get("satisfaction"):
            lines.append(
                '  "satisfaction_analysis": "Paragraph covering satisfaction scores, NPS if present, '
                'top positive themes, and one or two areas to watch. Quote specific scores '
                'and benchmark comparisons.",'
            )
        if sections.get("highlights"):
            lines.append('  "key_highlights": ["5 concise bullet strings, each a standalone insight"],')
        if sections.get("recommendations"):
            lines.append('  "recommendations": ["4 actionable bullet strings for the next event"],')

        lines.append(
            '  "closing_statement": "One professional closing paragraph thanking the client '
            'and expressing confidence in future events"'
        )
        lines.append("}")
        lines.append("")
        lines.append(
            "Rules: use only numbers present in the data — do not invent figures. "
            "If data for a section is absent, write a graceful note. "
            "Keep each paragraph 3-5 sentences. Professional, positive, client-facing tone."
        )

        return "\n".join(lines)

    # ── Claude call ────────────────────────────────────────────────────────────

    def _call_claude(self, user_prompt: str) -> dict:
        message = self.client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=2048,
            system=[
                {
                    "type": "text",
                    "text": SYSTEM_PROMPT,
                    "cache_control": {"type": "ephemeral"},
                }
            ],
            messages=[{"role": "user", "content": user_prompt}],
        )

        raw = message.content[0].text.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.MULTILINE)
        raw = re.sub(r"\s*```\s*$", "", raw, flags=re.MULTILINE)
        return json.loads(raw)

    # ── Public entry point ─────────────────────────────────────────────────────

    def generate(
        self,
        config: dict,
        attendance_summary: dict | None = None,
        survey_summary: dict | None = None,
        benchmark_context: dict | None = None,
        event_context: dict | None = None,
        analyst_findings: dict | None = None,
        memory_context: dict | None = None,
    ) -> bytes:
        prompt = self._build_user_prompt(
            config,
            attendance_summary,
            survey_summary,
            benchmark_context,
            event_context=event_context,
            analyst_findings=analyst_findings,
            memory_context=memory_context,
        )
        insights = self._call_claude(prompt)
        return self._build_pptx(config, insights, attendance_summary, survey_summary, benchmark_context)

    # ── Low-level drawing helpers ──────────────────────────────────────────────

    def _bg(self, slide, color: RGBColor):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _rect(self, slide, x, y, w, h, fill_color: RGBColor,
              border_color: RGBColor = None, border_pt: float = 0):
        shp = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
        if border_color:
            shp.line.color.rgb = border_color
            shp.line.width = Pt(border_pt)
        else:
            shp.line.fill.background()
        return shp

    def _txt(self, slide, text, x, y, w, h,
             size=14, bold=False, italic=False,
             color: RGBColor = MID, align=PP_ALIGN.LEFT,
             font=None, wrap=True, line_spacing=None):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font or TXT
        return box

    # ── Slide builders ─────────────────────────────────────────────────────────

    def _cover_slide(self, prs, config: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._bg(slide, NAVY)

        # Top indigo bar
        self._rect(slide, 0, 0, SW, 0.5, INDIGO)
        # Bottom indigo bar
        self._rect(slide, 0, SH - 0.55, SW, 0.55, INDIGO)

        # Event name
        self._txt(
            slide, config["event_name"].upper(),
            0.8, 1.75, SW - 1.6, 1.55,
            size=46, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, font=H1,
        )

        # Subtitle label
        self._txt(
            slide, "P O S T   E V E N T   R E P O R T",
            0.8, 3.45, SW - 1.6, 0.55,
            size=13, color=INDIGO,
            align=PP_ALIGN.CENTER, font=TXT,
        )

        # Meta
        meta = [str(config["event_date"])]
        if config.get("client_name"):
            meta.append(f"Prepared for: {config['client_name']}")
        if config.get("organizer_name"):
            meta.append(f"By: {config['organizer_name']}")

        self._txt(
            slide, "  ·  ".join(meta),
            0.8, 4.15, SW - 1.6, 0.5,
            size=11, color=LIGHT,
            align=PP_ALIGN.CENTER, font=TXT,
        )

    def _glance_slide(self, prs, config: dict,
                      attendance_summary: dict | None,
                      survey_summary: dict | None,
                      benchmark_context: dict | None):
        metrics: list[tuple[str, str]] = []

        if attendance_summary:
            reg = attendance_summary.get("registered", attendance_summary.get("total_rows"))
            if reg is not None:
                metrics.append(("Total Registered", f"{reg:,}"))
            if "checked_in" in attendance_summary:
                metrics.append(("Attended", f"{attendance_summary['checked_in']:,}"))
            if "attendance_rate_pct" in attendance_summary:
                metrics.append(("Show-up Rate", f"{attendance_summary['attendance_rate_pct']}%"))

        if survey_summary and survey_summary.get("rating_columns"):
            _, col_stats = next(iter(survey_summary["rating_columns"].items()))
            scale = "/10" if (col_stats.get("max") or 0) > 5 else "/5"
            metrics.append(("Avg Satisfaction", f"{col_stats['mean']}{scale}"))

        if not metrics:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header band
        self._rect(slide, 0, 0, SW, 0.62, NAVY)
        self._txt(slide, "AT A GLANCE",
                  0.45, 0.12, 8, 0.42,
                  size=19, bold=True, color=WHITE, font=H2)
        self._txt(slide, str(config["event_date"]),
                  SW - 2.6, 0.15, 2.25, 0.35,
                  size=11, color=LIGHT, align=PP_ALIGN.RIGHT)

        # Cards
        n = len(metrics)
        gap = 0.28
        card_w = (SW - 0.8 - gap * (n - 1)) / n
        card_h = 3.1
        card_top = 0.62 + (SH - 0.62 - card_h - 0.45) / 2

        for i, (label, value) in enumerate(metrics):
            cx = 0.4 + i * (card_w + gap)

            # Drop shadow (offset slightly)
            self._rect(slide, cx + 0.05, card_top + 0.05,
                       card_w, card_h, SHADOW)

            # Card body
            self._rect(slide, cx, card_top, card_w, card_h,
                       WHITE, BORDER, 0.5)

            # Top accent stripe
            self._rect(slide, cx, card_top, card_w, 0.09, INDIGO)

            # Value (big number)
            self._txt(slide, value,
                      cx + 0.1, card_top + 0.4, card_w - 0.2, 1.35,
                      size=44, bold=True, color=INDIGO,
                      align=PP_ALIGN.CENTER, font=H1)

            # Thin divider
            self._rect(slide, cx + 0.35, card_top + 1.9,
                       card_w - 0.7, 0.025, BORDER)

            # Label
            self._txt(slide, label,
                      cx + 0.1, card_top + 2.05, card_w - 0.2, 0.65,
                      size=12, color=MID,
                      align=PP_ALIGN.CENTER, font=TXT)

        # Benchmark attribution
        if benchmark_context and benchmark_context.get("source"):
            self._txt(slide,
                      f"Benchmarks sourced from {benchmark_context['source']}",
                      0, SH - 0.38, SW, 0.3,
                      size=9, color=LIGHT, align=PP_ALIGN.CENTER)

    def _content_slide(self, prs, heading: str,
                       body: str = None, bullets: list = None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header band
        self._rect(slide, 0, 0, SW, 0.65, INDIGO)
        self._txt(slide, heading.upper(),
                  0.45, 0.13, SW - 0.9, 0.44,
                  size=17, bold=True, color=WHITE, font=H2)

        # Footer line
        self._rect(slide, 0, SH - 0.38, SW, 0.025, BORDER)
        self._txt(slide, "CONFIDENTIAL",
                  0.45, SH - 0.36, 4, 0.3,
                  size=8, color=LIGHT)

        if body:
            box = slide.shapes.add_textbox(
                Inches(0.45), Inches(0.85),
                Inches(SW - 0.9), Inches(SH - 1.38),
            )
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.line_spacing = 1.55
            run = p.add_run()
            run.text = str(body)
            run.font.size = Pt(15)
            run.font.color.rgb = MID
            run.font.name = TXT

        if bullets:
            box = slide.shapes.add_textbox(
                Inches(0.45), Inches(0.85),
                Inches(SW - 0.9), Inches(SH - 1.38),
            )
            tf = box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(16)
                p.line_spacing = 1.3
                # Coloured bullet dot
                dot = p.add_run()
                dot.text = "●   "
                dot.font.color.rgb = INDIGO
                dot.font.size = Pt(9)
                dot.font.name = TXT
                # Bullet text
                txt = p.add_run()
                txt.text = str(item)
                txt.font.size = Pt(15)
                txt.font.color.rgb = MID
                txt.font.name = TXT

        return slide

    def _end_slide(self, prs, config: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._bg(slide, INDIGO)

        # Dark bottom stripe
        self._rect(slide, 0, SH - 0.5, SW, 0.5, INDIGO_D)

        # Event name
        self._txt(
            slide, config["event_name"].upper(),
            0.8, 2.0, SW - 1.6, 1.5,
            size=40, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, font=H1,
        )

        # Thank you line
        self._txt(
            slide, "Thank You",
            0.8, 3.65, SW - 1.6, 0.65,
            size=18, color=SOFT,
            align=PP_ALIGN.CENTER, font=H1,
        )

        # Confidential / meta
        meta = (
            f"Confidential  ·  {config.get('organizer_name', '')}  ·  "
            f"{date.today().strftime('%B %Y')}"
        )
        self._txt(
            slide, meta,
            0.8, 4.45, SW - 1.6, 0.5,
            size=10, color=SOFT,
            align=PP_ALIGN.CENTER, font=TXT,
        )

    # ── PPTX assembly ──────────────────────────────────────────────────────────

    def _build_pptx(
        self,
        config: dict,
        insights: dict,
        attendance_summary: dict | None,
        survey_summary: dict | None,
        benchmark_context: dict | None,
    ) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(SW)
        prs.slide_height = Inches(SH)

        self._cover_slide(prs, config)
        self._glance_slide(prs, config, attendance_summary, survey_summary, benchmark_context)

        section_map = [
            ("executive_summary",    "Executive Summary",                "paragraph"),
            ("attendance_analysis",  "Attendance & Reach",               "paragraph"),
            ("satisfaction_analysis","Attendee Satisfaction & Feedback",  "paragraph"),
            ("key_highlights",       "Key Highlights",                   "bullets"),
            ("recommendations",      "Recommendations for Future Events", "bullets"),
            ("closing_statement",    "Closing",                          "paragraph"),
        ]

        for key, heading, kind in section_map:
            if key not in insights:
                continue
            if kind == "paragraph":
                self._content_slide(prs, heading, body=str(insights[key]))
            else:
                self._content_slide(prs, heading, bullets=insights[key])

        self._end_slide(prs, config)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()
